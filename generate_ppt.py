#!/usr/bin/env python3
"""
Generate a premium PowerPoint presentation for the
Multi-Vendor Marketplace System – Final Year Internship Project.

Uses blue/purple gradient palette, Poppins-like typography,
glassmorphism-inspired cards, icons, and embedded screenshots.
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap

# ──────────────────────────────────────────────
# COLOUR PALETTE
# ──────────────────────────────────────────────
DEEP_BLUE     = RGBColor(0x1A, 0x1A, 0x4E)   # Dark navy
BRAND_BLUE    = RGBColor(0x3B, 0x5B, 0xDB)   # Primary blue
BRAND_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)   # Primary purple
LIGHT_PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)   # Accent
SOFT_BLUE     = RGBColor(0xD0, 0xEB, 0xFF)   # Card bg
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE     = RGBColor(0xF8, 0xF9, 0xFC)
DARK_TEXT      = RGBColor(0x1E, 0x1E, 0x2E)
BODY_TEXT      = RGBColor(0x4A, 0x4A, 0x5A)
MUTED_TEXT     = RGBColor(0x7A, 0x7A, 0x8C)
ACCENT_GREEN   = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED     = RGBColor(0xEF, 0x44, 0x44)
CARD_BG        = RGBColor(0xF0, 0xF0, 0xFF)   # Subtle lavender
GRADIENT_START = RGBColor(0x3B, 0x5B, 0xDB)
GRADIENT_END   = RGBColor(0x7C, 0x3A, 0xED)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

SCREENSHOT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "report-screenshots")

# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────

def _set_slide_bg(slide, color):
    """Set solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_gradient_bg(slide, color1, color2):
    """Set a two-stop gradient background (left to right)."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name="Poppins"):
    """Add a simple text box and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_paragraph(text_frame, text, font_size=14, bold=False,
                   color=BODY_TEXT, alignment=PP_ALIGN.LEFT,
                   font_name="Poppins", space_before=Pt(6),
                   space_after=Pt(2)):
    """Append a new paragraph to an existing text frame."""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def _add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG,
                      border_color=None, shadow=True):
    """Add a rounded rectangle (glassmorphism card)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()  # no border

    return shape


def _add_icon_circle(slide, left, top, size, fill_color, text="", emoji=""):
    """Add a coloured circle with optional text/emoji inside."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text or emoji:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = emoji or text
        p.font.size = Pt(int(size / Inches(1) * 14))
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Segoe UI Emoji"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def _add_gradient_bar(slide, left, top, width, height):
    """Add a decorative gradient bar across top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = BRAND_BLUE
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = BRAND_PURPLE
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()
    return shape


def _add_slide_number(slide, num, total):
    """Add slide number at bottom right."""
    _add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=MUTED_TEXT,
                 alignment=PP_ALIGN.RIGHT)


def _add_section_header(slide, title, subtitle=""):
    """Add a consistent section header with gradient underline."""
    _add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
                 title, font_size=32, bold=True, color=DEEP_BLUE)
    if subtitle:
        _add_textbox(slide, Inches(0.8), Inches(1.05), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=MUTED_TEXT)
    # Small accent bar under title
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0) if not subtitle else Inches(1.5),
        Inches(1.5), Inches(0.05)
    )
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = BRAND_BLUE
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = BRAND_PURPLE
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()


def _add_card_with_icon(slide, left, top, width, height, icon_emoji,
                        title, body_lines, icon_color=BRAND_BLUE):
    """Add a glassmorphism-style card with icon, title, and bullet points."""
    card = _add_rounded_rect(slide, left, top, width, height,
                             fill_color=WHITE, border_color=RGBColor(0xE0, 0xE0, 0xF0))

    # Icon circle
    icon_size = Inches(0.55)
    _add_icon_circle(slide, left + Inches(0.3), top + Inches(0.3),
                     icon_size, icon_color, emoji=icon_emoji)

    # Title
    _add_textbox(slide, left + Inches(1.0), top + Inches(0.3),
                 width - Inches(1.3), Inches(0.4),
                 title, font_size=16, bold=True, color=DEEP_BLUE)

    # Body bullets
    if body_lines:
        tf = _add_textbox(slide, left + Inches(0.4), top + Inches(0.85),
                          width - Inches(0.8), height - Inches(1.0),
                          "", font_size=12, color=BODY_TEXT)
        tf.paragraphs[0].text = ""
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = f"• {line}"
                tf.paragraphs[0].font.size = Pt(12)
                tf.paragraphs[0].font.color.rgb = BODY_TEXT
                tf.paragraphs[0].font.name = "Poppins"
            else:
                _add_paragraph(tf, f"• {line}", font_size=12, color=BODY_TEXT,
                              space_before=Pt(3), space_after=Pt(1))


def _add_feature_pill(slide, left, top, width, height, text,
                      fill_color=CARD_BG, text_color=BRAND_BLUE):
    """Add a rounded pill-shaped feature tag."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Poppins"
    p.alignment = PP_ALIGN.CENTER
    return shape


# ──────────────────────────────────────────────
# SLIDE BUILDERS
# ──────────────────────────────────────────────
TOTAL_SLIDES = 18

def slide_01_title(prs):
    """Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_gradient_bg(slide, DEEP_BLUE, RGBColor(0x2D, 0x1B, 0x69))

    # Decorative circles
    _add_icon_circle(slide, Inches(-0.5), Inches(-0.5), Inches(3),
                     RGBColor(0x4A, 0x3A, 0x8A))
    _add_icon_circle(slide, Inches(11), Inches(5.5), Inches(3),
                     RGBColor(0x3B, 0x2D, 0x7A))

    # Main content area
    card = _add_rounded_rect(slide, Inches(2.2), Inches(1.2), Inches(8.9),
                             Inches(5.1), fill_color=RGBColor(0xFF, 0xFF, 0xFF),
                             border_color=None)
    # Make card semi-transparent for glass effect via XML
    try:
        solidFill = card._element.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = srgb.makeelement(qn('a:alpha'), {'val': '90000'})
                srgb.append(alpha)
    except:
        pass

    # Top accent bar on card
    _add_gradient_bar(slide, Inches(2.2), Inches(1.2), Inches(8.9), Inches(0.08))

    # Title text
    _add_textbox(slide, Inches(2.8), Inches(1.6), Inches(7.7), Inches(0.7),
                 "🛒  MULTI-VENDOR MARKETPLACE SYSTEM", font_size=13,
                 bold=True, color=BRAND_PURPLE, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(2.8), Inches(2.2), Inches(7.7), Inches(1.0),
                 "Multi-Vendor\nMarketplace System", font_size=40,
                 bold=True, color=DEEP_BLUE, alignment=PP_ALIGN.CENTER)

    # Gradient line
    _add_gradient_bar(slide, Inches(5.3), Inches(3.3), Inches(2.7), Inches(0.05))

    _add_textbox(slide, Inches(2.8), Inches(3.6), Inches(7.7), Inches(0.5),
                 "Internship Final Year Project", font_size=18,
                 bold=False, color=BRAND_BLUE, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(2.8), Inches(4.3), Inches(7.7), Inches(0.4),
                 "Presented by  Rishiraj Sharma", font_size=16,
                 bold=True, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(2.8), Inches(4.75), Inches(7.7), Inches(0.4),
                 "Diploma in Computer Engineering", font_size=13,
                 color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(2.8), Inches(5.15), Inches(7.7), Inches(0.4),
                 "G. V. Acharya Polytechnic", font_size=13,
                 bold=True, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(2.8), Inches(5.55), Inches(7.7), Inches(0.4),
                 "Academic Year 2025 – 2026", font_size=12,
                 color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)


def slide_02_introduction(prs):
    """Introduction slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Introduction", "Understanding Multi-Vendor Marketplaces")
    _add_slide_number(slide, 2, TOTAL_SLIDES)

    # Definition card
    card = _add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(11.7),
                             Inches(1.3), fill_color=WHITE,
                             border_color=RGBColor(0xE0, 0xE0, 0xF0))
    _add_icon_circle(slide, Inches(1.1), Inches(1.85), Inches(0.55),
                     BRAND_BLUE, emoji="📖")
    tf = _add_textbox(slide, Inches(1.85), Inches(1.85), Inches(10.3), Inches(1.0),
                      "What is a Multi-Vendor Marketplace?", font_size=16,
                      bold=True, color=DEEP_BLUE)
    _add_paragraph(tf, "A multi-vendor marketplace is an e-commerce platform where multiple independent vendors "
                       "can list and sell their products through a single storefront, managed by a central admin.",
                   font_size=12, color=BODY_TEXT)

    # Real-world examples — 3 cards
    examples = [
        ("🛍️", "Amazon", "Global leader with\nmillions of vendors", BRAND_BLUE),
        ("📦", "Flipkart", "India's largest\nonline marketplace", BRAND_PURPLE),
        ("🎨", "Etsy", "Handmade & vintage\nspecialty marketplace", ACCENT_GREEN),
    ]
    x_start = Inches(0.8)
    for i, (icon, title, desc, color) in enumerate(examples):
        left = x_start + Inches(i * 4.05)
        card = _add_rounded_rect(slide, left, Inches(3.3), Inches(3.7),
                                 Inches(1.7), fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.25), Inches(3.5),
                         Inches(0.5), color, emoji=icon)
        _add_textbox(slide, left + Inches(0.9), Inches(3.45), Inches(2.5),
                     Inches(0.35), title, font_size=16, bold=True, color=DEEP_BLUE)
        _add_textbox(slide, left + Inches(0.9), Inches(3.85), Inches(2.5),
                     Inches(0.7), desc, font_size=12, color=BODY_TEXT)

    # Benefits section
    _add_textbox(slide, Inches(0.8), Inches(5.3), Inches(3), Inches(0.4),
                 "Key Benefits", font_size=18, bold=True, color=DEEP_BLUE)

    benefits = [
        ("✅", "Wider product selection for customers"),
        ("✅", "Low startup cost for vendors"),
        ("✅", "Revenue through commissions for admin"),
        ("✅", "Centralized order & inventory management"),
    ]
    for i, (icon, text) in enumerate(benefits):
        _add_textbox(slide, Inches(0.8), Inches(5.8 + i * 0.38), Inches(12),
                     Inches(0.35), f"{icon}  {text}", font_size=13,
                     color=BODY_TEXT)


def slide_03_problem(prs):
    """Problem Statement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Problem Statement", "Challenges that inspired this project")
    _add_slide_number(slide, 3, TOTAL_SLIDES)

    problems = [
        ("🏪", "Limited Online Presence",
         "Small businesses lack the technical resources to build their own e-commerce platforms.",
         ACCENT_RED),
        ("📊", "No Centralized Management",
         "Managing vendors, products, and orders across multiple systems is inefficient.",
         ACCENT_ORANGE),
        ("👥", "Customer Reach",
         "Individual vendors struggle to attract traffic and compete with large retailers.",
         BRAND_PURPLE),
        ("🔒", "Security Concerns",
         "Lack of secure authentication and payment handling puts users at risk.",
         BRAND_BLUE),
    ]

    for i, (icon, title, desc, color) in enumerate(problems):
        col = i % 2
        row = i // 2
        left = Inches(0.8) + Inches(col * 6.2)
        top = Inches(1.8) + Inches(row * 2.6)

        card = _add_rounded_rect(slide, left, top, Inches(5.7), Inches(2.2),
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.3), top + Inches(0.35),
                         Inches(0.6), color, emoji=icon)
        _add_textbox(slide, left + Inches(1.1), top + Inches(0.35),
                     Inches(4.0), Inches(0.4), title, font_size=18,
                     bold=True, color=DEEP_BLUE)
        _add_textbox(slide, left + Inches(0.3), top + Inches(1.1),
                     Inches(5.0), Inches(0.9), desc, font_size=13,
                     color=BODY_TEXT)


def slide_04_objectives(prs):
    """Project Objectives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Project Objectives", "Goals driving the development")
    _add_slide_number(slide, 4, TOTAL_SLIDES)

    objectives = [
        ("🎯", "Centralized Marketplace", "Build a unified platform for vendors and customers"),
        ("🔐", "Secure Authentication", "Implement role-based login with password hashing"),
        ("👤", "Vendor Management", "Admin-approved vendor registration & store management"),
        ("📦", "Product Management", "Full CRUD with categories, images, inventory tracking"),
        ("🚚", "Order Tracking", "Real-time order status updates for all stakeholders"),
        ("📱", "Responsive Design", "Mobile-first, works on all screen sizes"),
    ]

    for i, (icon, title, desc) in enumerate(objectives):
        col = i % 3
        row = i // 3
        left = Inches(0.8) + Inches(col * 4.1)
        top = Inches(1.7) + Inches(row * 2.8)

        card = _add_rounded_rect(slide, left, top, Inches(3.7), Inches(2.4),
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(1.4), top + Inches(0.25),
                         Inches(0.65), BRAND_BLUE if i % 2 == 0 else BRAND_PURPLE,
                         emoji=icon)
        _add_textbox(slide, left + Inches(0.2), top + Inches(1.05),
                     Inches(3.3), Inches(0.4), title, font_size=15,
                     bold=True, color=DEEP_BLUE, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, left + Inches(0.2), top + Inches(1.5),
                     Inches(3.3), Inches(0.7), desc, font_size=12,
                     color=BODY_TEXT, alignment=PP_ALIGN.CENTER)


def slide_05_users(prs):
    """System Users."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "System Users", "Three key roles in the marketplace")
    _add_slide_number(slide, 5, TOTAL_SLIDES)

    users = [
        ("👤", "Customer", BRAND_BLUE, [
            "Browse & search products",
            "Add to cart & wishlist",
            "Place & track orders",
            "Write product reviews",
            "Manage profile & addresses",
        ]),
        ("🏪", "Vendor", BRAND_PURPLE, [
            "Register & manage store",
            "Add / edit / delete products",
            "Manage inventory & stock",
            "Process incoming orders",
            "View sales analytics",
        ]),
        ("🛡️", "Admin", ACCENT_GREEN, [
            "Approve / reject vendors",
            "Manage all users & products",
            "View platform analytics",
            "Monitor orders & payments",
            "System-wide oversight",
        ]),
    ]

    for i, (icon, title, color, items) in enumerate(users):
        left = Inches(0.6) + Inches(i * 4.2)
        top = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.3)

        card = _add_rounded_rect(slide, left, top, w, h,
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))

        # Top gradient accent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.07)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        _add_icon_circle(slide, left + Inches(1.4), top + Inches(0.3),
                         Inches(0.8), color, emoji=icon)

        _add_textbox(slide, left + Inches(0.2), top + Inches(1.25), Inches(3.4),
                     Inches(0.4), title, font_size=22, bold=True,
                     color=DEEP_BLUE, alignment=PP_ALIGN.CENTER)

        for j, item in enumerate(items):
            _add_textbox(slide, left + Inches(0.4), top + Inches(1.85 + j * 0.55),
                         Inches(3.0), Inches(0.4),
                         f"▸  {item}", font_size=12, color=BODY_TEXT)


def slide_06_architecture(prs):
    """System Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "System Architecture", "High-level system overview")
    _add_slide_number(slide, 6, TOTAL_SLIDES)

    # ── Presentation Layer ──
    _add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
                 "PRESENTATION LAYER", font_size=11, bold=True,
                 color=BRAND_BLUE)

    pres_items = [("👤", "Customer\nBrowser", BRAND_BLUE),
                  ("🏪", "Vendor\nPanel", BRAND_PURPLE),
                  ("🛡️", "Admin\nDashboard", ACCENT_GREEN)]
    for i, (icon, label, color) in enumerate(pres_items):
        left = Inches(1.5) + Inches(i * 3.8)
        card = _add_rounded_rect(slide, left, Inches(2.0), Inches(2.8),
                                 Inches(1.2), fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.15), Inches(2.15),
                         Inches(0.45), color, emoji=icon)
        _add_textbox(slide, left + Inches(0.7), Inches(2.1), Inches(1.9),
                     Inches(1.0), label, font_size=13, bold=True,
                     color=DEEP_BLUE, alignment=PP_ALIGN.CENTER)

    # Arrow down
    _add_textbox(slide, Inches(6.0), Inches(3.25), Inches(1.5), Inches(0.4),
                 "▼  ▼  ▼", font_size=20, color=BRAND_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # ── Application Layer ──
    _add_textbox(slide, Inches(0.8), Inches(3.7), Inches(11), Inches(0.3),
                 "APPLICATION LAYER", font_size=11, bold=True,
                 color=BRAND_PURPLE)

    app_card = _add_rounded_rect(slide, Inches(1.5), Inches(4.05), Inches(10.3),
                                 Inches(1.1), fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))

    app_items = ["PHP Backend", "REST API", "Session Mgmt", "Auth & RBAC"]
    for i, item in enumerate(app_items):
        _add_feature_pill(slide, Inches(1.8) + Inches(i * 2.6), Inches(4.25),
                          Inches(2.2), Inches(0.6), f"⚙️ {item}",
                          fill_color=RGBColor(0xED, 0xEF, 0xFF),
                          text_color=BRAND_PURPLE)

    # Arrow down
    _add_textbox(slide, Inches(6.0), Inches(5.2), Inches(1.5), Inches(0.4),
                 "▼  ▼  ▼", font_size=20, color=BRAND_PURPLE,
                 alignment=PP_ALIGN.CENTER)

    # ── Data Layer ──
    _add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.3),
                 "DATA LAYER", font_size=11, bold=True, color=ACCENT_GREEN)

    db_card = _add_rounded_rect(slide, Inches(3.0), Inches(5.9), Inches(7.3),
                                Inches(1.1), fill_color=WHITE,
                                border_color=RGBColor(0xE0, 0xE0, 0xF0))

    db_items = ["MySQL Database", "File Storage"]
    for i, item in enumerate(db_items):
        _add_feature_pill(slide, Inches(3.5) + Inches(i * 3.5), Inches(6.1),
                          Inches(2.8), Inches(0.6), f"🗄️ {item}",
                          fill_color=RGBColor(0xE8, 0xFA, 0xF0),
                          text_color=ACCENT_GREEN)


def slide_07_techstack(prs):
    """Technology Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Technology Stack", "Tools & technologies powering the system")
    _add_slide_number(slide, 7, TOTAL_SLIDES)

    categories = [
        ("Frontend", "🎨", [
            ("HTML5", "Structure & semantics"),
            ("CSS3", "Styling & animations"),
            ("JavaScript", "Client-side logic"),
            ("Bootstrap 5", "Responsive framework"),
        ], BRAND_BLUE),
        ("Backend", "⚙️", [
            ("PHP", "Server-side logic"),
            ("MySQL", "Relational database"),
            ("XAMPP", "Local dev server"),
        ], BRAND_PURPLE),
        ("Dev Tools", "🛠️", [
            ("VS Code", "Code editor"),
            ("Git & GitHub", "Version control"),
            ("Chrome DevTools", "Debugging"),
        ], ACCENT_GREEN),
    ]

    for ci, (cat_name, cat_icon, techs, color) in enumerate(categories):
        left = Inches(0.6) + Inches(ci * 4.2)
        top = Inches(1.7)
        w = Inches(3.8)
        h = Inches(5.3)

        card = _add_rounded_rect(slide, left, top, w, h,
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        # Accent top
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.06)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        _add_icon_circle(slide, left + Inches(1.4), top + Inches(0.2),
                         Inches(0.7), color, emoji=cat_icon)
        _add_textbox(slide, left + Inches(0.2), top + Inches(1.05), Inches(3.4),
                     Inches(0.4), cat_name, font_size=18, bold=True,
                     color=DEEP_BLUE, alignment=PP_ALIGN.CENTER)

        for j, (tech, desc) in enumerate(techs):
            ty = top + Inches(1.65 + j * 0.85)
            _add_rounded_rect(slide, left + Inches(0.25), ty, Inches(3.3),
                              Inches(0.7), fill_color=RGBColor(0xF5, 0xF5, 0xFF),
                              border_color=None)
            _add_textbox(slide, left + Inches(0.4), ty + Inches(0.05),
                         Inches(3.0), Inches(0.3), tech, font_size=13,
                         bold=True, color=DEEP_BLUE)
            _add_textbox(slide, left + Inches(0.4), ty + Inches(0.35),
                         Inches(3.0), Inches(0.3), desc, font_size=10,
                         color=MUTED_TEXT)


def slide_08_database(prs):
    """Database Design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Database Design", "16 interconnected tables powering the platform")
    _add_slide_number(slide, 8, TOTAL_SLIDES)

    tables = [
        ("👤", "Users", "id, name, email, role,\npassword_hash, status"),
        ("🏪", "Vendors", "store_name, owner, gstin,\nstatus, commission_rate"),
        ("📦", "Products", "title, price, stock, sku,\ncategory_id, vendor_id"),
        ("📂", "Categories", "name, slug, icon,\nparent_id, image"),
        ("🛒", "Cart / Items", "user_id, product_id,\nquantity, price"),
        ("❤️", "Wishlist", "user_id, product_id,\ncreated_at"),
        ("📋", "Orders", "order_number, status,\ntotal_amount, address"),
        ("💳", "Payments", "transaction_id, method,\namount, status"),
        ("⭐", "Reviews", "rating, comment,\nvendor_response"),
    ]

    for i, (icon, name, fields) in enumerate(tables):
        col = i % 3
        row = i // 3
        left = Inches(0.6) + Inches(col * 4.2)
        top = Inches(1.7) + Inches(row * 1.85)
        w = Inches(3.8)
        h = Inches(1.6)

        card = _add_rounded_rect(slide, left, top, w, h,
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.15), top + Inches(0.15),
                         Inches(0.45), BRAND_BLUE if col == 0 else (BRAND_PURPLE if col == 1 else ACCENT_GREEN),
                         emoji=icon)
        _add_textbox(slide, left + Inches(0.7), top + Inches(0.15),
                     Inches(2.9), Inches(0.35), name, font_size=15,
                     bold=True, color=DEEP_BLUE)
        _add_textbox(slide, left + Inches(0.2), top + Inches(0.65),
                     Inches(3.4), Inches(0.8), fields, font_size=10,
                     color=BODY_TEXT)


def slide_09_customer_module(prs):
    """Customer Module."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Customer Module", "End-to-end shopping experience")
    _add_slide_number(slide, 9, TOTAL_SLIDES)

    # Flow steps
    steps = [
        ("1", "Register", "Create account with\nemail verification", BRAND_BLUE),
        ("2", "Login", "Secure session-based\nauthentication", BRAND_BLUE),
        ("3", "Browse", "Search, filter &\nexplore products", BRAND_PURPLE),
        ("4", "Wishlist", "Save favourite\nproducts for later", BRAND_PURPLE),
        ("5", "Cart", "Add items, adjust\nquantities & total", ACCENT_GREEN),
        ("6", "Checkout", "Enter shipping info\n& place order", ACCENT_GREEN),
        ("7", "Track", "View order history\n& live status", RGBColor(0xF5, 0x9E, 0x0B)),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        left = Inches(0.3) + Inches(i * 1.82)
        top = Inches(1.7)
        w = Inches(1.65)
        h = Inches(2.8)

        card = _add_rounded_rect(slide, left, top, w, h,
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        # Number circle
        _add_icon_circle(slide, left + Inches(0.45), top + Inches(0.2),
                         Inches(0.55), color, text=num)
        _add_textbox(slide, left + Inches(0.1), top + Inches(0.9),
                     Inches(1.45), Inches(0.3), title, font_size=13,
                     bold=True, color=DEEP_BLUE, alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, left + Inches(0.1), top + Inches(1.3),
                     Inches(1.45), Inches(1.0), desc, font_size=10,
                     color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

        # Arrow
        if i < len(steps) - 1:
            _add_textbox(slide, left + Inches(1.55), top + Inches(1.0),
                         Inches(0.4), Inches(0.4), "→", font_size=20,
                         color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

    # Bottom highlights
    highlights = [
        "🔍 Advanced search with filters",
        "📱 Fully responsive mobile layout",
        "⭐ Product reviews & ratings",
        "📜 Detailed order history"
    ]
    for i, h in enumerate(highlights):
        _add_feature_pill(slide, Inches(0.5) + Inches(i * 3.15), Inches(5.0),
                          Inches(3.0), Inches(0.55), h,
                          fill_color=RGBColor(0xED, 0xEF, 0xFF),
                          text_color=BRAND_PURPLE)


def slide_10_vendor_module(prs):
    """Vendor Module."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Vendor Module", "Complete store management toolkit")
    _add_slide_number(slide, 10, TOTAL_SLIDES)

    features = [
        ("📋", "Registration", "Apply as vendor with store\ndetails, await admin approval", BRAND_BLUE),
        ("📊", "Dashboard", "Overview of sales, orders,\nrevenue & analytics", BRAND_PURPLE),
        ("📦", "Product CRUD", "Create, read, update & delete\nproducts with images", ACCENT_GREEN),
        ("📈", "Inventory", "Track stock levels, low-stock\nalerts, bulk updates", ACCENT_ORANGE),
        ("🚚", "Orders", "View & process incoming\norders, update status", BRAND_BLUE),
        ("💰", "Analytics", "Sales reports, revenue trends,\ntop-selling products", BRAND_PURPLE),
    ]

    for i, (icon, title, desc, color) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.6) + Inches(col * 4.2)
        top = Inches(1.7) + Inches(row * 2.7)

        card = _add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.3),
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.25), top + Inches(0.3),
                         Inches(0.6), color, emoji=icon)
        _add_textbox(slide, left + Inches(1.05), top + Inches(0.3),
                     Inches(2.5), Inches(0.35), title, font_size=17,
                     bold=True, color=DEEP_BLUE)
        _add_textbox(slide, left + Inches(0.25), top + Inches(1.1),
                     Inches(3.3), Inches(1.0), desc, font_size=12,
                     color=BODY_TEXT)


def slide_11_admin_module(prs):
    """Admin Module."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Admin Module", "Platform control & oversight")
    _add_slide_number(slide, 11, TOTAL_SLIDES)

    features = [
        ("👥", "User Management", "View, activate, suspend\nor delete user accounts"),
        ("✅", "Vendor Approval", "Review vendor applications,\napprove or reject"),
        ("📦", "Product Oversight", "Monitor all products,\nremove violations"),
        ("📊", "Reports & Analytics", "Revenue, orders, users\nplatform-wide stats"),
        ("🛡️", "Dashboard", "Real-time overview of\nall platform activity"),
        ("📋", "Activity Logs", "Track all actions for\naudit & compliance"),
    ]

    colors = [BRAND_BLUE, ACCENT_GREEN, BRAND_PURPLE, ACCENT_ORANGE, BRAND_BLUE, BRAND_PURPLE]
    for i, (icon, title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.6) + Inches(col * 4.2)
        top = Inches(1.7) + Inches(row * 2.7)

        card = _add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.3),
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.25), top + Inches(0.3),
                         Inches(0.6), colors[i], emoji=icon)
        _add_textbox(slide, left + Inches(1.05), top + Inches(0.3),
                     Inches(2.5), Inches(0.35), title, font_size=17,
                     bold=True, color=DEEP_BLUE)
        _add_textbox(slide, left + Inches(0.25), top + Inches(1.1),
                     Inches(3.3), Inches(1.0), desc, font_size=12,
                     color=BODY_TEXT)


def slide_12_workflow(prs):
    """Workflow / Flowchart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Application Workflow", "End-to-end user journey flowchart")
    _add_slide_number(slide, 12, TOTAL_SLIDES)

    steps = [
        ("1", "Register /\nLogin", BRAND_BLUE),
        ("2", "Browse\nProducts", BRAND_BLUE),
        ("3", "Add to\nCart", BRAND_PURPLE),
        ("4", "Checkout\n& Pay", BRAND_PURPLE),
        ("5", "Order\nConfirmed", ACCENT_GREEN),
        ("6", "Vendor\nProcesses", ACCENT_ORANGE),
        ("7", "Shipped &\nDelivered", ACCENT_GREEN),
        ("8", "Admin\nMonitors", RGBColor(0x4A, 0x3A, 0x8A)),
    ]

    # Draw as connected horizontal flow
    y_center = Inches(3.2)
    step_w = Inches(1.3)
    step_h = Inches(1.3)
    gap = Inches(0.25)
    total_w = len(steps) * (step_w + gap) - gap
    x_start = (SLIDE_WIDTH - total_w) / 2

    for i, (num, label, color) in enumerate(steps):
        left = x_start + i * (step_w + gap)
        top = y_center

        # Step shape (rounded rect)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, step_w, step_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Poppins"
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(10)
        p2.font.color.rgb = WHITE
        p2.font.name = "Poppins"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)

        # Arrow
        if i < len(steps) - 1:
            _add_textbox(slide, left + step_w - Inches(0.05), y_center + Inches(0.45),
                         gap + Inches(0.1), Inches(0.4),
                         "→", font_size=22, color=MUTED_TEXT,
                         alignment=PP_ALIGN.CENTER)

    # Descriptions below
    desc_row = [
        ("Customer Actions", "Register → Browse → Cart → Checkout", BRAND_BLUE),
        ("System Actions", "Order Confirmed → Payment Processed", BRAND_PURPLE),
        ("Vendor Actions", "Receive Order → Process → Ship", ACCENT_GREEN),
        ("Admin Actions", "Monitor → Report → Manage", RGBColor(0x4A, 0x3A, 0x8A)),
    ]

    for i, (title, desc, color) in enumerate(desc_row):
        left = Inches(0.5) + Inches(i * 3.2)
        top = Inches(5.2)
        _add_rounded_rect(slide, left, top, Inches(3.0), Inches(1.5),
                          fill_color=WHITE, border_color=RGBColor(0xE0, 0xE0, 0xF0))
        # Accent left bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, Inches(0.07), Inches(1.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        _add_textbox(slide, left + Inches(0.2), top + Inches(0.15),
                     Inches(2.6), Inches(0.3), title, font_size=12,
                     bold=True, color=DEEP_BLUE)
        _add_textbox(slide, left + Inches(0.2), top + Inches(0.55),
                     Inches(2.6), Inches(0.8), desc, font_size=10,
                     color=BODY_TEXT)


def slide_13_features(prs):
    """Key Features."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Key Features", "What makes this platform stand out")
    _add_slide_number(slide, 13, TOTAL_SLIDES)

    features = [
        ("🏪", "Multi-Vendor\nSupport"),
        ("📱", "Responsive\nDesign"),
        ("🔐", "Secure\nAuthentication"),
        ("🔍", "Search &\nFilter"),
        ("❤️", "Wishlist"),
        ("🛒", "Shopping\nCart"),
        ("📦", "Inventory\nManagement"),
        ("🚚", "Order\nTracking"),
        ("📊", "Sales\nDashboard"),
    ]

    colors = [BRAND_BLUE, BRAND_PURPLE, ACCENT_GREEN, ACCENT_ORANGE,
              ACCENT_RED, BRAND_BLUE, BRAND_PURPLE, ACCENT_GREEN, ACCENT_ORANGE]

    for i, (icon, label) in enumerate(features):
        col = i % 3
        row = i // 3
        left = Inches(0.8) + Inches(col * 4.1)
        top = Inches(1.7) + Inches(row * 1.85)

        card = _add_rounded_rect(slide, left, top, Inches(3.7), Inches(1.6),
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.25), top + Inches(0.35),
                         Inches(0.65), colors[i], emoji=icon)
        _add_textbox(slide, left + Inches(1.1), top + Inches(0.3),
                     Inches(2.3), Inches(1.0), label, font_size=16,
                     bold=True, color=DEEP_BLUE)


def slide_14_advantages(prs):
    """Advantages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Advantages", "Why choose this marketplace platform")
    _add_slide_number(slide, 14, TOTAL_SLIDES)

    advantages = [
        ("✨", "Easy to Use", "Intuitive interface designed for all user types — customers, vendors, and admins.", BRAND_BLUE),
        ("🏪", "Helps Small Businesses", "Provides a ready-made platform for small businesses to reach customers without building their own.", BRAND_PURPLE),
        ("😊", "Better Customer Experience", "Smooth browsing, wishlists, reviews, and streamlined checkout process.", ACCENT_GREEN),
        ("📈", "Scalable Architecture", "Modular design supports growing number of vendors, products, and users.", ACCENT_ORANGE),
        ("🔒", "Secure & Reliable", "Role-based access control, password hashing, and input validation.", BRAND_BLUE),
        ("🎯", "Centralized Management", "Single admin dashboard to monitor and control the entire platform.", BRAND_PURPLE),
    ]

    for i, (icon, title, desc, color) in enumerate(advantages):
        col = i % 2
        row = i // 2
        left = Inches(0.8) + Inches(col * 6.2)
        top = Inches(1.7) + Inches(row * 1.85)

        card = _add_rounded_rect(slide, left, top, Inches(5.7), Inches(1.6),
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.25), top + Inches(0.35),
                         Inches(0.6), color, emoji=icon)
        _add_textbox(slide, left + Inches(1.05), top + Inches(0.25),
                     Inches(4.2), Inches(0.35), title, font_size=16,
                     bold=True, color=DEEP_BLUE)
        _add_textbox(slide, left + Inches(1.05), top + Inches(0.7),
                     Inches(4.2), Inches(0.7), desc, font_size=11,
                     color=BODY_TEXT)


def slide_15_future_scope(prs):
    """Future Scope."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Future Scope", "Planned enhancements & roadmap")
    _add_slide_number(slide, 15, TOTAL_SLIDES)

    # Timeline-style layout
    items = [
        ("💳", "Payment Gateway", "Integrate Razorpay / Stripe\nfor online payments", "Phase 1"),
        ("📱", "Mobile App", "Native Android/iOS app\nfor on-the-go shopping", "Phase 1"),
        ("🤖", "AI Recommendations", "ML-powered product suggestions\nbased on user behavior", "Phase 2"),
        ("🚛", "Delivery Tracking", "Real-time delivery tracking\nwith map integration", "Phase 2"),
        ("🎫", "Coupons & Offers", "Discount codes, flash sales,\nand promotional campaigns", "Phase 3"),
        ("📧", "Email Notifications", "Automated emails for orders,\nshipping & promotions", "Phase 3"),
        ("🌐", "Multi-language", "Support for Hindi, Marathi,\nand other languages", "Phase 3"),
    ]

    for i, (icon, title, desc, phase) in enumerate(items):
        col = i % 2
        row = i // 2 if i < 6 else 2
        if i == 6:
            col = 0
        left = Inches(0.8) + Inches(col * 6.2)
        top = Inches(1.7) + Inches(row * 1.4)

        if i == 6:  # Center the last item
            left = Inches(3.8)

        card = _add_rounded_rect(slide, left, top, Inches(5.7), Inches(1.15),
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))
        _add_icon_circle(slide, left + Inches(0.2), top + Inches(0.2),
                         Inches(0.5), BRAND_BLUE if "1" in phase else (BRAND_PURPLE if "2" in phase else ACCENT_GREEN),
                         emoji=icon)
        _add_textbox(slide, left + Inches(0.85), top + Inches(0.12),
                     Inches(3.5), Inches(0.3), title, font_size=14,
                     bold=True, color=DEEP_BLUE)
        _add_textbox(slide, left + Inches(0.85), top + Inches(0.45),
                     Inches(3.5), Inches(0.6), desc, font_size=10,
                     color=BODY_TEXT)

        # Phase badge
        _add_feature_pill(slide, left + Inches(4.5), top + Inches(0.3),
                          Inches(1.0), Inches(0.4), phase,
                          fill_color=RGBColor(0xED, 0xEF, 0xFF),
                          text_color=BRAND_PURPLE)


def slide_16_screenshots(prs):
    """Screenshots Section with embedded project screenshots."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Screenshots", "Live application screenshots")
    _add_slide_number(slide, 16, TOTAL_SLIDES)

    screenshots = [
        ("homepage_hero_1786027775121.png", "Home Page"),
        ("shop_page_1786027831752.png", "Shop Page"),
        ("product_detail_1786027932910.png", "Product Detail"),
        ("customer_dashboard_1786028045135.png", "Customer Dashboard"),
        ("vendor_dashboard_1786028132409.png", "Vendor Dashboard"),
        ("admin_dashboard_1786028190010.png", "Admin Dashboard"),
    ]

    for i, (filename, label) in enumerate(screenshots):
        col = i % 3
        row = i // 3
        left = Inches(0.6) + Inches(col * 4.2)
        top = Inches(1.7) + Inches(row * 2.7)
        w = Inches(3.8)
        h = Inches(2.3)

        # Card background
        card = _add_rounded_rect(slide, left, top, w, h,
                                 fill_color=WHITE,
                                 border_color=RGBColor(0xE0, 0xE0, 0xF0))

        # Try to add screenshot
        img_path = os.path.join(SCREENSHOT_DIR, filename)
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(
                    img_path, left + Inches(0.1), top + Inches(0.1),
                    w - Inches(0.2), h - Inches(0.5)
                )
            except Exception as e:
                _add_textbox(slide, left + Inches(0.3), top + Inches(0.6),
                             Inches(3.2), Inches(0.4),
                             f"[Screenshot: {label}]", font_size=12,
                             color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)
        else:
            _add_textbox(slide, left + Inches(0.3), top + Inches(0.6),
                         Inches(3.2), Inches(0.4),
                         f"[Screenshot: {label}]", font_size=12,
                         color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

        # Label
        _add_textbox(slide, left, top + h - Inches(0.35), w, Inches(0.3),
                     label, font_size=11, bold=True, color=DEEP_BLUE,
                     alignment=PP_ALIGN.CENTER)


def slide_17_summary(prs):
    """Project Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, OFF_WHITE)
    _add_section_header(slide, "Project Summary", "Key takeaways & achievements")
    _add_slide_number(slide, 17, TOTAL_SLIDES)

    sections = [
        ("🏆", "Key Achievements", [
            "Built a full-stack multi-vendor marketplace",
            "Implemented 3 distinct user roles with RBAC",
            "Designed 16-table relational database schema",
            "Created responsive UI with Bootstrap 5",
        ], BRAND_BLUE),
        ("💡", "Skills Learned", [
            "Full-stack web development (PHP + MySQL)",
            "Database design & normalization",
            "UI/UX design principles",
            "Version control with Git & GitHub",
        ], BRAND_PURPLE),
        ("⚙️", "Technologies Used", [
            "HTML5, CSS3, JavaScript, Bootstrap 5",
            "PHP 8.x, MySQL / MariaDB",
            "XAMPP, VS Code, Chrome DevTools",
            "Git & GitHub for collaboration",
        ], ACCENT_GREEN),
        ("🎓", "Internship Outcome", [
            "Successfully completed the project",
            "Gained real-world development experience",
            "Learned project management & planning",
            "Ready for deployment & future scaling",
        ], ACCENT_ORANGE),
    ]

    for i, (icon, title, items, color) in enumerate(sections):
        col = i % 2
        row = i // 2
        left = Inches(0.8) + Inches(col * 6.2)
        top = Inches(1.7) + Inches(row * 2.8)

        _add_card_with_icon(slide, left, top, Inches(5.7), Inches(2.5),
                            icon, title, items, icon_color=color)


def slide_18_thankyou(prs):
    """Thank You Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, DEEP_BLUE, RGBColor(0x2D, 0x1B, 0x69))

    # Decorative circles
    _add_icon_circle(slide, Inches(10.5), Inches(-1), Inches(3.5),
                     RGBColor(0x4A, 0x3A, 0x8A))
    _add_icon_circle(slide, Inches(-1), Inches(5), Inches(3),
                     RGBColor(0x3B, 0x2D, 0x7A))

    # Glass card
    card = _add_rounded_rect(slide, Inches(3.0), Inches(1.5), Inches(7.3),
                             Inches(4.5), fill_color=RGBColor(0xFF, 0xFF, 0xFF))
    try:
        solidFill = card._element.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha = srgb.makeelement(qn('a:alpha'), {'val': '88000'})
                srgb.append(alpha)
    except:
        pass

    _add_gradient_bar(slide, Inches(3.0), Inches(1.5), Inches(7.3), Inches(0.08))

    _add_textbox(slide, Inches(3.5), Inches(2.0), Inches(6.3), Inches(1.0),
                 "Thank You!", font_size=48, bold=True, color=DEEP_BLUE,
                 alignment=PP_ALIGN.CENTER)

    _add_gradient_bar(slide, Inches(5.5), Inches(3.1), Inches(2.3), Inches(0.05))

    _add_textbox(slide, Inches(3.5), Inches(3.4), Inches(6.3), Inches(0.6),
                 "Questions & Answers", font_size=22, color=BRAND_PURPLE,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(3.5), Inches(4.1), Inches(6.3), Inches(0.4),
                 "Rishiraj Sharma", font_size=16, bold=True, color=DARK_TEXT,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(3.5), Inches(4.5), Inches(6.3), Inches(0.4),
                 "Diploma in Computer Engineering  •  G. V. Acharya Polytechnic",
                 font_size=12, color=BODY_TEXT, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(3.5), Inches(5.0), Inches(6.3), Inches(0.4),
                 "Multi-Vendor Marketplace System  |  2025–2026",
                 font_size=11, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────

def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = Emu(12192000)   # 13.333 inches
    prs.slide_height = Emu(6858000)   # 7.5 inches

    builders = [
        slide_01_title,
        slide_02_introduction,
        slide_03_problem,
        slide_04_objectives,
        slide_05_users,
        slide_06_architecture,
        slide_07_techstack,
        slide_08_database,
        slide_09_customer_module,
        slide_10_vendor_module,
        slide_11_admin_module,
        slide_12_workflow,
        slide_13_features,
        slide_14_advantages,
        slide_15_future_scope,
        slide_16_screenshots,
        slide_17_summary,
        slide_18_thankyou,
    ]

    for builder in builders:
        print(f"  Building: {builder.__doc__}")
        builder(prs)

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                            "Multi_Vendor_Marketplace_PPT.pptx")
    prs.save(out_path)
    print(f"\n[OK] Presentation saved to: {out_path}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
